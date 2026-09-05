#!/usr/bin/env python3
"""Generate and validate the editable PowerPoint presentation of the TCC.

The PTC presentation is used only as a narrative and visual reference.  This
generator intentionally reads the current versioned visualizations and keeps
the scientific caveats visible in the deck.  The committed PPTX is the
click-to-open artifact; this script makes its regeneration auditable.

Usage from the repository root::

    python scripts/generate_tcc_presentation.py
    python scripts/generate_tcc_presentation.py --check
"""

from __future__ import annotations

import argparse
import hashlib
import sys
from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = ROOT / "results" / "tcc" / "presentation" / "ensino_personalizado_de_matematica_tcc.pptx"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
ORANGE = RGBColor(234, 88, 12)
VIOLET = RGBColor(124, 58, 237)
INDIGO = RGBColor(79, 70, 229)
RED = RGBColor(185, 28, 28)
WHITE = RGBColor(255, 255, 255)
PALE_BLUE = RGBColor(239, 246, 255)
PALE_TEAL = RGBColor(240, 253, 250)
PALE_GREEN = RGBColor(240, 253, 244)
PALE_AMBER = RGBColor(255, 251, 235)
PALE_VIOLET = RGBColor(245, 243, 255)


EXPECTED_TITLES = [
    "Ensino Personalizado de Matemática",
    "O problema de pesquisa",
    "Questão e objetivo",
    "Objetivos específicos",
    "A base conceitual: quatro níveis de interpretação",
    "Desenho metodológico",
    "Do registro bruto à população retida",
    "Fluxo PRISMA do snapshot",
    "Deduplicação: o que foi confirmado",
    "Panorama descritivo do snapshot",
    "Distribuição temporal e fontes",
    "População retida e síntese empírica",
    "Apreciação metodológica pelo MMAT 2018",
    "O que a síntese sustenta",
    "Lacunas documentadas",
    "Da evidência à especificação",
    "Especificação conceitual do protótipo",
    "Contribuições e limites",
    "Obrigado",
]

CURRENT_MARKERS = (
    "11.904",
    "11.877",
    "9.391",
    "2.486",
    "2.468",
    "18",
    "27",
    "2015–2026",
    "MMAT 2018",
)

# These values are present in historical PTC material and must not leak into
# the current TCC deck as if they were current evidence.
FORBIDDEN_HISTORICAL_MARKERS = (
    "9.431",
    "2.517",
    "6.914",
    "1.883",
    "85%",
    "0% BNCC",
    "35%",
    "17 estudos incluídos",
    "17 estudos incluidos",
    "6915",
    "6916",
    "6917",
    "6918",
    "6919",
    "6920",
    "6921",
    "6923",
)


def rgb(value: RGBColor) -> RGBColor:
    return value


def add_shape(slide, shape_type, x, y, width, height, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(shape_type, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_text(
    slide,
    text: str,
    x,
    y,
    width,
    height,
    *,
    size=18,
    color=SLATE,
    bold=False,
    font="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.06,
    name=None,
):
    box = slide.shapes.add_textbox(x, y, width, height)
    if name:
        box.name = name
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_lines(
    slide,
    lines: Sequence[tuple[str, str, RGBColor, bool]],
    x,
    y,
    width,
    height,
    *,
    size=16,
    color=SLATE,
    bullet=False,
    line_spacing=1.12,
):
    box = slide.shapes.add_textbox(x, y, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)
    for index, (label, body, label_color, label_bold) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(6)
        paragraph.line_spacing = line_spacing
        if bullet:
            paragraph.text = "• "
            paragraph.runs[0].font.name = "Aptos"
            paragraph.runs[0].font.size = Pt(size)
            paragraph.runs[0].font.color.rgb = label_color
        first = paragraph.add_run()
        first.text = label
        first.font.name = "Aptos"
        first.font.size = Pt(size)
        first.font.bold = label_bold
        first.font.color.rgb = label_color
        if body:
            second = paragraph.add_run()
            second.text = body
            second.font.name = "Aptos"
            second.font.size = Pt(size)
            second.font.color.rgb = color
    return box


def add_bullets(slide, items: Sequence[str], x, y, width, height, *, size=17, color=SLATE, spacing=8):
    box = slide.shapes.add_textbox(x, y, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.level = 0
        paragraph.space_after = Pt(spacing)
        paragraph.line_spacing = 1.08
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
    return box


def add_title(slide, title: str, number: int):
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), BLUE)
    add_text(
        slide,
        title,
        Inches(0.55),
        Inches(0.35),
        Inches(11.7),
        Inches(0.55),
        size=27,
        color=NAVY,
        bold=True,
        name="tcc-title",
    )
    add_text(
        slide,
        f"TCC · Ensino Personalizado de Matemática  |  {number:02d}",
        Inches(0.58),
        Inches(7.12),
        Inches(8.5),
        Inches(0.2),
        size=8.5,
        color=MUTED,
    )
    add_text(
        slide,
        "IFC · Ciência da Computação · 2026",
        Inches(9.35),
        Inches(7.12),
        Inches(3.4),
        Inches(0.2),
        size=8.5,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_callout(slide, title: str, body: str, x, y, width, height, *, fill=PALE_BLUE, accent=BLUE, title_size=16, body_size=13.5):
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height, fill)
    card.line.color.rgb = accent
    card.line.width = Pt(1.25)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(0.09), height, accent)
    add_text(slide, title, x + Inches(0.18), y + Inches(0.1), width - Inches(0.3), Inches(0.35), size=title_size, color=accent, bold=True)
    add_text(slide, body, x + Inches(0.18), y + Inches(0.48), width - Inches(0.3), height - Inches(0.55), size=body_size, color=SLATE)
    return card


def add_metric(slide, value: str, label: str, x, y, width, height, *, fill=PALE_BLUE, accent=BLUE, value_size=25, label_size=11.5):
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height, fill)
    card.line.color.rgb = accent
    card.line.width = Pt(1)
    add_text(slide, value, x, y + Inches(0.16), width, Inches(0.48), size=value_size, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + Inches(0.1), y + Inches(0.68), width - Inches(0.2), height - Inches(0.72), size=label_size, color=SLATE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return card


def add_image_contain(slide, path: Path, x, y, width, height):
    if not path.exists():
        raise FileNotFoundError(f"Visualization not found: {path}")
    with Image.open(path) as image:
        image_width, image_height = image.size
    image_ratio = image_width / image_height
    box_ratio = width / height
    if image_ratio >= box_ratio:
        picture_width = width
        picture_height = int(width / image_ratio)
        picture_x = x
        picture_y = y + int((height - picture_height) / 2)
    else:
        picture_height = height
        picture_width = int(height * image_ratio)
        picture_x = x + int((width - picture_width) / 2)
        picture_y = y
    picture = slide.shapes.add_picture(str(path), picture_x, picture_y, width=picture_width, height=picture_height)
    # python-pptx stores the media part as image.png inside the ZIP package.
    # Preserve the source basename on the shape so the validator can verify
    # which canonical visualization was embedded.
    picture.name = path.name
    return picture


def new_slide(prs: Presentation, title: str, number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_title(slide, title, number)
    return slide


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    # 1 — Cover.  The narrative follows the PTC structure while making the
    # current TCC scope explicit.
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), SLIDE_HEIGHT, BLUE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(6.8), Inches(13.13), Inches(0.7), TEAL)
    add_text(slide, "TRABALHO DE CONCLUSÃO DE CURSO", Inches(0.82), Inches(0.7), Inches(6), Inches(0.25), size=12, color=TEAL, bold=True)
    add_text(slide, "Ensino Personalizado\nde Matemática", Inches(0.8), Inches(1.35), Inches(8.3), Inches(1.45), size=36, color=WHITE, bold=True, name="tcc-title")
    add_text(slide, "Oportunidades e Técnicas Computacionais", Inches(0.84), Inches(3.05), Inches(8.4), Inches(0.45), size=22, color=RGBColor(186, 230, 253), bold=True)
    add_text(slide, "Revisão sistemática da literatura e especificação conceitual de protótipo", Inches(0.84), Inches(3.7), Inches(7.5), Inches(0.62), size=17, color=RGBColor(226, 232, 240))
    add_text(slide, "Thales Ferreira Batista\nProf. Dr. Rafael Zanin · Orientador\nProf. Dr. Manassés Ribeiro · Coorientador", Inches(0.84), Inches(5.1), Inches(6.6), Inches(0.9), size=14, color=WHITE)
    add_text(slide, "IFC · Videira\nSnapshot adjudicado: 03/09/2026 · recorte 2015–2026", Inches(8.95), Inches(5.25), Inches(3.35), Inches(0.62), size=11.5, color=RGBColor(226, 232, 240), align=PP_ALIGN.RIGHT)

    # 2 — Problem.
    slide = new_slide(prs, EXPECTED_TITLES[1], 2)
    add_text(slide, "Turmas heterogêneas produzem muitas evidências, mas a evidência só se torna útil quando é interpretada no contexto pedagógico.", Inches(0.75), Inches(1.25), Inches(11.85), Inches(0.65), size=21, color=NAVY, bold=True)
    add_callout(slide, "Contexto", "Conhecimentos prévios, ritmos e dificuldades variam entre estudantes e ao longo do tempo.", Inches(0.75), Inches(2.2), Inches(3.75), Inches(1.7), fill=PALE_BLUE, accent=BLUE)
    add_callout(slide, "Oportunidade", "Técnicas computacionais podem organizar registros, revelar padrões e apoiar o acompanhamento.", Inches(4.8), Inches(2.2), Inches(3.75), Inches(1.7), fill=PALE_TEAL, accent=TEAL)
    add_callout(slide, "Responsabilidade", "A saída computacional não substitui a avaliação do professor nem representa, sozinha, a aprendizagem.", Inches(8.85), Inches(2.2), Inches(3.75), Inches(1.7), fill=PALE_AMBER, accent=AMBER)
    add_text(slide, "O trabalho investiga como a literatura articula técnicas, evidências educacionais e decisões de personalização — e transforma esse conhecimento em uma especificação auditável.", Inches(1.15), Inches(5.0), Inches(11.0), Inches(0.8), size=20, color=SLATE, align=PP_ALIGN.CENTER)

    # 3 — Question and objective.
    slide = new_slide(prs, EXPECTED_TITLES[2], 3)
    add_callout(slide, "Problema de pesquisa", "Como identificar e sintetizar, por meio de revisão sistemática, as principais técnicas computacionais aplicadas ao ensino de matemática e converter essas evidências em uma especificação de protótipo que apoie o professor?", Inches(0.75), Inches(1.25), Inches(7.2), Inches(2.1), fill=PALE_BLUE, accent=BLUE, title_size=17, body_size=16)
    add_callout(slide, "Objetivo geral", "Mapear e analisar sistematicamente as aplicações e elaborar uma especificação técnica e pedagógica para apoiar a interpretação de evidências sobre competências.", Inches(8.25), Inches(1.25), Inches(4.3), Inches(2.1), fill=PALE_TEAL, accent=TEAL, title_size=17, body_size=16)
    add_text(slide, "Perguntas orientadoras", Inches(0.78), Inches(4.15), Inches(4.2), Inches(0.35), size=18, color=NAVY, bold=True)
    add_bullets(slide, [
        "Quais técnicas são aplicadas à educação matemática?",
        "Como essas aplicações são avaliadas nos contextos estudados?",
        "Quais lacunas e limitações aparecem na literatura?",
        "Que requisitos podem orientar uma ferramenta de apoio?",
    ], Inches(0.82), Inches(4.55), Inches(11.2), Inches(1.55), size=16)

    # 4 — Objectives.
    slide = new_slide(prs, EXPECTED_TITLES[3], 4)
    left = [
        ("OE1", "Revisão sistemática, relatada com apoio do PRISMA 2020, sobre estudos de 2015–2026."),
        ("OE2", "Identificação e categorização das abordagens computacionais."),
        ("OE3", "Classificação das finalidades pedagógicas das aplicações."),
        ("OE4", "Análise crítica das metodologias e limitações dos estudos."),
    ]
    right = [
        ("OE5", "Mapeamento de lacunas técnicas, pedagógicas, metodológicas e éticas."),
        ("OE6", "Manutenção de um pipeline automatizado e auditável para coleta, processamento e exportação."),
        ("OE7", "Derivação de requisitos, critérios de dados e modelos, protocolo de avaliação e arquitetura de referência."),
    ]
    for index, (code, body) in enumerate(left):
        add_callout(slide, code, body, Inches(0.75), Inches(1.25 + index * 1.25), Inches(5.65), Inches(0.95), fill=PALE_BLUE, accent=BLUE, title_size=15, body_size=13.5)
    for index, (code, body) in enumerate(right):
        add_callout(slide, code, body, Inches(6.9), Inches(1.25 + index * 1.55), Inches(5.65), Inches(1.25), fill=PALE_TEAL, accent=TEAL, title_size=15, body_size=13.5)
    add_text(slide, "A apresentação distingue objetivos executados de atividades que permanecem fora do escopo, como implementação funcional e validação com participantes.", Inches(1.0), Inches(6.45), Inches(11.25), Inches(0.42), size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 5 — Theoretical foundation.
    slide = new_slide(prs, EXPECTED_TITLES[4], 5)
    add_text(slide, "A revisão pode reconhecer padrões; ela não observa diretamente todos os processos cognitivos e sociais envolvidos na aprendizagem.", Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.55), size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    cards = [
        ("Desempenho\nobservado", "Registro de uma tarefa: acerto, nota, tentativa, estratégia ou tempo.", BLUE, PALE_BLUE),
        ("Proficiência\nestimada", "Inferência sobre uma escala de conhecimentos e habilidades.", TEAL, PALE_TEAL),
        ("Competência", "Mobilização integrada de conhecimentos, procedimentos, estratégias e atitudes.", VIOLET, PALE_VIOLET),
        ("Aprendizagem", "Transformação construída ao longo do tempo, com compreensão, autonomia e transferência.", AMBER, PALE_AMBER),
    ]
    for index, (title, body, accent, fill) in enumerate(cards):
        x = Inches(0.72 + index * 3.12)
        add_callout(slide, title, body, x, Inches(2.05), Inches(2.75), Inches(2.25), fill=fill, accent=accent, title_size=17, body_size=13)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.9), NAVY)
    add_text(slide, "Princípio de interpretação: a saída computacional apoia a análise docente; não constitui diagnóstico definitivo.", Inches(1.75), Inches(5.25), Inches(9.8), Inches(0.35), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 6 — Method.
    slide = new_slide(prs, EXPECTED_TITLES[5], 6)
    method_cards = [
        ("PRISMA 2020", "Diretriz de relato transparente da revisão.", BLUE, PALE_BLUE),
        ("PICOS", "População, intervenção, comparação, resultados e desenho.", TEAL, PALE_TEAL),
        ("4 fontes", "Semantic Scholar, OpenAlex, Crossref e CORE.", VIOLET, PALE_VIOLET),
        ("Escopo", "Inglês/português · 2015–2026 · corte em 31/08/2026.", AMBER, PALE_AMBER),
    ]
    for index, (title, body, accent, fill) in enumerate(method_cards):
        add_callout(slide, title, body, Inches(0.72 + (index % 2) * 6.15), Inches(1.2 + (index // 2) * 1.35), Inches(5.55), Inches(1.05), fill=fill, accent=accent, title_size=15, body_size=13.5)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.45), Inches(11.8), Inches(1.35), PALE_BLUE)
    add_text(slide, "Estratégia canônica", Inches(1.0), Inches(4.7), Inches(2.5), Inches(0.3), size=16, color=BLUE, bold=True)
    add_text(slide, "72 consultas compostas (48 em inglês + 24 em português), combinando domínio matemático, técnica computacional e contexto educacional.", Inches(1.0), Inches(5.05), Inches(10.9), Inches(0.45), size=17, color=SLATE)
    add_text(slide, "O número descreve a composição versionada da estratégia; não é uma contagem retrospectiva de chamadas HTTP concluídas, pois não há log histórico completo.", Inches(1.0), Inches(5.58), Inches(10.9), Inches(0.35), size=11.5, color=MUTED)

    # 7 — Flow metrics.
    slide = new_slide(prs, EXPECTED_TITLES[6], 7)
    add_text(slide, "Contagens do snapshot adjudicado", Inches(0.75), Inches(1.05), Inches(5), Inches(0.3), size=18, color=NAVY, bold=True)
    metrics = [
        ("11.904", "registros identificados", BLUE, PALE_BLUE),
        ("27", "remoções por identidade\n25 DOI + 2 URL", TEAL, PALE_TEAL),
        ("11.877", "registros na triagem", GREEN, PALE_GREEN),
        ("9.391", "excluídos na triagem", AMBER, PALE_AMBER),
        ("2.486", "na elegibilidade", ORANGE, RGBColor(255, 247, 237)),
        ("2.468", "excluídos na elegibilidade", VIOLET, PALE_VIOLET),
        ("18", "retidos operacionalmente", INDIGO, RGBColor(238, 242, 255)),
    ]
    positions = [(0.75, 1.6), (3.15, 1.6), (5.55, 1.6), (7.95, 1.6), (1.95, 4.15), (4.35, 4.15), (6.75, 4.15)]
    for (value, label, accent, fill), (x, y) in zip(metrics, positions):
        add_metric(slide, value, label, Inches(x), Inches(y), Inches(2.0), Inches(1.5), fill=fill, accent=accent, value_size=24, label_size=10.5)
    add_text(slide, "A remoção determinística por DOI/URL precede a triagem. Igualdade de título permaneceu como candidato à auditoria semântica.", Inches(1.1), Inches(6.2), Inches(11.0), Inches(0.35), size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 8 — PRISMA image.
    slide = new_slide(prs, EXPECTED_TITLES[7], 8)
    add_image_contain(slide, ROOT / "research" / "exports" / "visualizations" / "prisma_flow.png", Inches(0.9), Inches(1.05), Inches(11.55), Inches(5.7))
    add_text(slide, "Figura versionada e sincronizada com os artefatos públicos do pipeline.", Inches(1.0), Inches(6.75), Inches(11.25), Inches(0.25), size=11.5, color=MUTED, align=PP_ALIGN.CENTER)

    # 9 — Deduplication.
    slide = new_slide(prs, EXPECTED_TITLES[8], 9)
    add_image_contain(slide, ROOT / "research" / "exports" / "visualizations" / "selection_funnel.png", Inches(0.7), Inches(1.15), Inches(6.15), Inches(5.4))
    add_callout(slide, "Confirmadas no fluxo", "27 linhas excedentes por identidade bibliográfica determinística: 25 DOI normalizado e 2 URL exata.", Inches(7.15), Inches(1.35), Inches(5.2), Inches(1.25), fill=PALE_TEAL, accent=TEAL, title_size=16, body_size=14)
    add_callout(slide, "Mantidas para auditoria", "232 excedentes apenas por título normalizado permaneceram como candidatos a revisão semântica.", Inches(7.15), Inches(2.95), Inches(5.2), Inches(1.25), fill=PALE_AMBER, accent=AMBER, title_size=16, body_size=14)
    add_callout(slide, "Leitura científica", "Títulos semelhantes podem ser versões, erratas ou obras distintas. Por isso, não foram tratados automaticamente como duplicatas confirmadas.", Inches(7.15), Inches(4.55), Inches(5.2), Inches(1.35), fill=PALE_BLUE, accent=BLUE, title_size=16, body_size=13.5)

    # 10 — Descriptive snapshot.
    slide = new_slide(prs, EXPECTED_TITLES[9], 10)
    add_image_contain(slide, ROOT / "research" / "exports" / "visualizations" / "techniques_distribution.png", Inches(0.55), Inches(1.3), Inches(8.0), Inches(5.15))
    add_text(slide, "Frequências mais recorrentes", Inches(8.85), Inches(1.35), Inches(3.5), Inches(0.3), size=17, color=NAVY, bold=True)
    add_metric(slide, "6.399", "técnica não especificada", Inches(8.9), Inches(1.85), Inches(3.3), Inches(0.95), fill=PALE_BLUE, accent=BLUE, value_size=21, label_size=10.5)
    add_metric(slide, "1.073", "assessment", Inches(8.9), Inches(2.95), Inches(3.3), Inches(0.95), fill=PALE_TEAL, accent=TEAL, value_size=21, label_size=10.5)
    add_metric(slide, "863", "IA / inteligência artificial", Inches(8.9), Inches(4.05), Inches(3.3), Inches(0.95), fill=PALE_VIOLET, accent=VIOLET, value_size=21, label_size=10.5)
    add_metric(slide, "771", "machine learning", Inches(8.9), Inches(5.15), Inches(3.3), Inches(0.95), fill=PALE_AMBER, accent=AMBER, value_size=21, label_size=10.5)
    add_text(slide, "As categorias podem se sobrepor e são calculadas sobre os 11.877 registros após a remoção determinística. Não representam qualidade ou eficácia.", Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.3), size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # 11 — Time and databases.
    slide = new_slide(prs, EXPECTED_TITLES[10], 11)
    add_image_contain(slide, ROOT / "research" / "exports" / "visualizations" / "papers_by_year.png", Inches(0.6), Inches(1.25), Inches(6.0), Inches(2.7))
    add_image_contain(slide, ROOT / "research" / "exports" / "visualizations" / "database_coverage.png", Inches(6.75), Inches(1.25), Inches(6.0), Inches(2.7))
    add_callout(slide, "Período", "O snapshot foi delimitado a 2015–2026, com data de corte em 31 de agosto de 2026.", Inches(0.85), Inches(4.35), Inches(5.65), Inches(1.15), fill=PALE_BLUE, accent=BLUE, title_size=15, body_size=13.5)
    add_callout(slide, "Fontes", "As quatro fontes foram complementares; os totais brutos incluem registros posteriormente reconhecidos como duplicatas.", Inches(6.95), Inches(4.35), Inches(5.55), Inches(1.15), fill=PALE_TEAL, accent=TEAL, title_size=15, body_size=13.5)
    add_text(slide, "Visualizações descritivas: não inferem representatividade, qualidade ou efeito pedagógico.", Inches(1.0), Inches(6.35), Inches(11.2), Inches(0.3), size=13, color=MUTED, align=PP_ALIGN.CENTER)

    # 12 — Retained population.
    slide = new_slide(prs, EXPECTED_TITLES[11], 12)
    add_metric(slide, "18", "registros retidos operacionalmente", Inches(0.85), Inches(1.35), Inches(3.2), Inches(1.35), fill=PALE_BLUE, accent=INDIGO, value_size=29, label_size=12)
    add_metric(slide, "17", "candidatos empíricos provisórios", Inches(4.95), Inches(1.35), Inches(3.2), Inches(1.35), fill=PALE_TEAL, accent=TEAL, value_size=29, label_size=12)
    add_metric(slide, "1", "protocolo ou proposta contextual", Inches(9.05), Inches(1.35), Inches(3.2), Inches(1.35), fill=PALE_AMBER, accent=AMBER, value_size=29, label_size=12)
    add_callout(slide, "O que entra na síntese", "A síntese de evidências e a apreciação MMAT aplicável consideram os 17 candidatos empíricos. O protocolo/proposta contextual foi mantido para rastreabilidade do mapeamento, mas não sustenta resultado empírico concluído.", Inches(0.85), Inches(3.35), Inches(5.65), Inches(1.85), fill=PALE_BLUE, accent=BLUE, title_size=17, body_size=14.5)
    add_callout(slide, "Padrão observado", "Predomínio de aplicações voltadas à predição de desempenho e à estimativa de proficiência, com recorrência de modelos supervisionados como Random Forest, SVM, árvores de decisão e redes neurais.", Inches(6.85), Inches(3.35), Inches(5.65), Inches(1.85), fill=PALE_TEAL, accent=TEAL, title_size=17, body_size=14.5)
    add_text(slide, "Os resultados de artigos individuais permanecem condicionados à população, ao instrumento, à métrica e ao desenho de cada estudo.", Inches(1.0), Inches(6.35), Inches(11.2), Inches(0.32), size=13.5, color=MUTED, align=PP_ALIGN.CENTER)

    # 13 — MMAT.
    slide = new_slide(prs, EXPECTED_TITLES[12], 13)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.25), Inches(3.2), Inches(4.55), NAVY)
    add_text(slide, "MMAT 2018", Inches(1.05), Inches(1.65), Inches(2.7), Inches(0.45), size=27, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "avaliação por critério\nsem nota global", Inches(1.05), Inches(2.35), Inches(2.7), Inches(0.85), size=19, color=RGBColor(186, 230, 253), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Instrumento aplicado conforme o desenho metodológico de cada estudo empírico.", Inches(1.1), Inches(4.15), Inches(2.6), Inches(0.75), size=13.5, color=RGBColor(226, 232, 240), align=PP_ALIGN.CENTER)
    add_callout(slide, "Como foi registrado", "Cinco critérios por desenho, com respostas Sim, Não ou Não é possível determinar. As limitações permanecem visíveis por estudo.", Inches(4.55), Inches(1.25), Inches(3.85), Inches(1.55), fill=PALE_BLUE, accent=BLUE, title_size=16, body_size=14)
    add_callout(slide, "Estado atual", "Apreciação preliminar por um único revisor: nove registros com texto primário revisado e oito com base em resumo/metadados.", Inches(8.8), Inches(1.25), Inches(3.75), Inches(1.55), fill=PALE_AMBER, accent=AMBER, title_size=16, body_size=14)
    add_callout(slide, "O que não foi feito", "Não há média, ranking ou categoria agregada de qualidade. A adjudicação, os localizadores e a recuperação de fontes ainda precisam ser consolidados.", Inches(4.55), Inches(3.35), Inches(8.0), Inches(1.55), fill=PALE_TEAL, accent=TEAL, title_size=16, body_size=14)
    add_text(slide, "O MMAT organiza a leitura das limitações; não transforma uma apreciação preliminar em classificação definitiva.", Inches(1.0), Inches(6.35), Inches(11.2), Inches(0.32), size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 14 — Findings and boundaries.
    slide = new_slide(prs, EXPECTED_TITLES[13], 14)
    add_callout(slide, "A síntese sustenta", "A literatura oferece um mapa auditável de técnicas, finalidades e formas de avaliação aplicadas à educação matemática.", Inches(0.8), Inches(1.35), Inches(5.75), Inches(1.55), fill=PALE_GREEN, accent=GREEN, title_size=17, body_size=15)
    add_callout(slide, "A heterogeneidade importa", "Populações, instrumentos, variáveis e métricas diferentes impedem comparar diretamente todos os resultados ou combiná-los em uma síntese quantitativa.", Inches(6.8), Inches(1.35), Inches(5.75), Inches(1.55), fill=PALE_BLUE, accent=BLUE, title_size=17, body_size=15)
    add_callout(slide, "A síntese não sustenta", "Uma técnica não pode ser declarada superior em geral, nem uma acurácia de um estudo ser transferida automaticamente para outra escola ou população.", Inches(0.8), Inches(3.65), Inches(5.75), Inches(1.55), fill=PALE_AMBER, accent=AMBER, title_size=17, body_size=15)
    add_callout(slide, "Implicação", "Avaliação técnica, explicabilidade, contexto curricular e interpretação docente precisam acompanhar qualquer uso futuro.", Inches(6.8), Inches(3.65), Inches(5.75), Inches(1.55), fill=PALE_VIOLET, accent=VIOLET, title_size=17, body_size=15)
    add_text(slide, "Desempenho reportado em um artigo não equivale a evidência geral de eficácia pedagógica.", Inches(1.0), Inches(6.35), Inches(11.2), Inches(0.32), size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 15 — Gaps.
    slide = new_slide(prs, EXPECTED_TITLES[14], 15)
    gaps = [
        ("Explicabilidade", "Pouca clareza sobre como as saídas chegam à interpretação pedagógica.", BLUE, PALE_BLUE),
        ("Currículo", "Alinhamento explícito com referenciais curriculares aparece de forma limitada.", TEAL, PALE_TEAL),
        ("Participação docente", "Decisões de projeto nem sempre incorporam o professor como intérprete.", VIOLET, PALE_VIOLET),
        ("Equidade", "Análises de possíveis vieses e diferenças entre grupos são escassas.", AMBER, PALE_AMBER),
        ("Reprodutibilidade", "Dados, códigos e detalhes de preparação nem sempre estão disponíveis.", ORANGE, RGBColor(255, 247, 237)),
        ("Contexto", "Validações permanecem restritas a populações e cenários específicos.", INDIGO, RGBColor(238, 242, 255)),
    ]
    for index, (title, body, accent, fill) in enumerate(gaps):
        x = Inches(0.75 + (index % 3) * 4.1)
        y = Inches(1.25 + (index // 3) * 2.0)
        add_callout(slide, title, body, x, y, Inches(3.65), Inches(1.55), fill=fill, accent=accent, title_size=15, body_size=13.5)
    add_text(slide, "Essas lacunas foram tratadas como requisitos de projeto — não como prova de que uma solução futura será eficaz.", Inches(1.0), Inches(6.35), Inches(11.2), Inches(0.32), size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 16 — From evidence to specification.
    slide = new_slide(prs, EXPECTED_TITLES[15], 16)
    stages = [
        ("Evidências", "revisão + fundamentação", BLUE, PALE_BLUE),
        ("Lacunas", "técnicas, pedagógicas, metodológicas e éticas", TEAL, PALE_TEAL),
        ("Requisitos", "funcionais e não funcionais", VIOLET, PALE_VIOLET),
        ("Protocolo", "dados, modelos, métricas e limites", AMBER, PALE_AMBER),
        ("Arquitetura", "referência conceitual auditável", INDIGO, RGBColor(238, 242, 255)),
    ]
    for index, (title, body, accent, fill) in enumerate(stages):
        x = Inches(0.55 + index * 2.55)
        add_callout(slide, title, body, x, Inches(2.05), Inches(2.2), Inches(2.0), fill=fill, accent=accent, title_size=15, body_size=12.5)
        if index < len(stages) - 1:
            add_text(slide, "→", Inches(2.78 + index * 2.55), Inches(2.72), Inches(0.38), Inches(0.35), size=24, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.35), Inches(5.05), Inches(10.55), Inches(0.8), NAVY)
    add_text(slide, "O resultado é uma especificação conceitual: documentada, verificável e coerente com o escopo executado.", Inches(1.65), Inches(5.3), Inches(9.95), Inches(0.3), size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 17 — Prototype specification.
    slide = new_slide(prs, EXPECTED_TITLES[16], 17)
    add_image_contain(slide, ROOT / "results" / "tcc" / "images" / "arquitetura_prototipo.png", Inches(0.65), Inches(1.15), Inches(6.35), Inches(5.45))
    add_text(slide, "Componentes previstos", Inches(7.35), Inches(1.25), Inches(4.4), Inches(0.35), size=18, color=NAVY, bold=True)
    add_bullets(slide, [
        "ingestão e validação das fontes",
        "preparação e documentação das variáveis",
        "modelagem e comparação de alternativas",
        "avaliação, explicabilidade e análise de erros",
        "apresentação para interpretação docente",
    ], Inches(7.35), Inches(1.75), Inches(5.0), Inches(2.45), size=15.5)
    add_callout(slide, "Limite do artefato", "Não é uma aplicação funcional: não houve treinamento com uma base definitiva, métricas próprias ou validação com participantes.", Inches(7.35), Inches(4.65), Inches(5.0), Inches(1.35), fill=PALE_AMBER, accent=AMBER, title_size=16, body_size=14)

    # 18 — Contributions and limitations.
    slide = new_slide(prs, EXPECTED_TITLES[17], 18)
    add_callout(slide, "Contribuições", "• síntese estruturada da literatura\n• pipeline automatizado e versionado\n• separação entre relevância e avaliação metodológica\n• especificação técnica e pedagógica auditável", Inches(0.8), Inches(1.25), Inches(5.75), Inches(3.1), fill=PALE_GREEN, accent=GREEN, title_size=18, body_size=15)
    add_callout(slide, "Limites assumidos", "• um único pesquisador\n• parte da apreciação baseada em resumo/metadados\n• fontes e idiomas delimitados\n• MMAT preliminar\n• sem aplicação funcional ou eficácia pedagógica própria", Inches(6.8), Inches(1.25), Inches(5.75), Inches(3.1), fill=PALE_AMBER, accent=AMBER, title_size=18, body_size=15)
    add_text(slide, "Conclusão", Inches(0.85), Inches(5.05), Inches(2.0), Inches(0.35), size=19, color=NAVY, bold=True)
    add_text(slide, "O trabalho respondeu ao escopo ao articular revisão sistemática, fundamentação pedagógica, apreciação metodológica e especificação conceitual — mantendo o professor no centro da interpretação.", Inches(0.85), Inches(5.55), Inches(11.45), Inches(0.7), size=19, color=SLATE, bold=True, align=PP_ALIGN.CENTER)

    # 19 — Closing.
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), TEAL)
    add_text(slide, "Obrigado", Inches(0.7), Inches(1.75), Inches(11.9), Inches(0.8), size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER, name="tcc-title")
    add_text(slide, "Perguntas e discussão", Inches(1.0), Inches(2.8), Inches(11.3), Inches(0.45), size=23, color=RGBColor(186, 230, 253), align=PP_ALIGN.CENTER)
    add_text(slide, "Ensino Personalizado de Matemática\nOportunidades e Técnicas Computacionais", Inches(1.2), Inches(4.35), Inches(10.9), Inches(0.8), size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Thales Ferreira Batista · Ciência da Computação · IFC Videira\nSnapshot adjudicado: 03/09/2026 · recorte temporal: 2015–2026", Inches(1.2), Inches(5.8), Inches(10.9), Inches(0.55), size=12.5, color=RGBColor(226, 232, 240), align=PP_ALIGN.CENTER)

    return prs


def slide_text(slide) -> str:
    values: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            values.append(shape.text)
    return "\n".join(values)


def validate_presentation(path: Path) -> list[str]:
    errors: list[str] = []
    if not path.exists():
        return [f"PPTX ausente: {path}"]
    try:
        prs = Presentation(str(path))
    except Exception as exc:  # pragma: no cover - defensive validation path
        return [f"PPTX inválido: {exc}"]

    if len(prs.slides) != len(EXPECTED_TITLES):
        errors.append(f"número de slides inesperado: {len(prs.slides)} (esperado {len(EXPECTED_TITLES)})")

    all_text = "\n".join(slide_text(slide) for slide in prs.slides)
    for marker in CURRENT_MARKERS:
        if marker not in all_text:
            errors.append(f"marcador científico atual ausente: {marker}")
    for marker in FORBIDDEN_HISTORICAL_MARKERS:
        if marker in all_text:
            errors.append(f"marcador histórico/operacional proibido encontrado: {marker}")

    for index, expected_title in enumerate(EXPECTED_TITLES):
        if index >= len(prs.slides):
            break
        title_shapes = [shape for shape in prs.slides[index].shapes if shape.name == "tcc-title"]
        title_text = title_shapes[0].text.replace("\n", " ") if title_shapes else ""
        normalized_expected = expected_title.replace("\n", " ")
        if normalized_expected not in title_text:
            errors.append(f"título do slide {index + 1} inesperado: {title_text!r}")

    expected_images = {
        "prisma_flow.png",
        "selection_funnel.png",
        "techniques_distribution.png",
        "papers_by_year.png",
        "database_coverage.png",
        "arquitetura_prototipo.png",
    }
    expected_image_sources = {
        name: ROOT / "research" / "exports" / "visualizations" / name
        for name in expected_images
        if name != "arquitetura_prototipo.png"
    }
    expected_image_sources["arquitetura_prototipo.png"] = ROOT / "results" / "tcc" / "images" / "arquitetura_prototipo.png"
    embedded_images: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.shape_type == 13:  # picture
                continue
            embedded_images.add(shape.name)
            source = expected_image_sources.get(shape.name)
            if source is not None and source.exists():
                source_hash = hashlib.sha256(source.read_bytes()).hexdigest()
                embedded_hash = hashlib.sha256(shape.image.blob).hexdigest()
                if source_hash != embedded_hash:
                    errors.append(f"imagem incorporada desatualizada: {shape.name}")
    missing_images = expected_images - embedded_images
    if missing_images:
        errors.append(f"imagens esperadas não incorporadas: {', '.join(sorted(missing_images))}")
    return errors


def generate(output: Path) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.core_properties.title = "Ensino Personalizado de Matemática — Apresentação do TCC"
    prs.core_properties.subject = "Revisão sistemática da literatura e especificação conceitual de protótipo"
    prs.core_properties.author = "Thales Ferreira Batista"
    prs.core_properties.comments = "Gerado por scripts/generate_tcc_presentation.py a partir dos artefatos versionados do TCC."
    prs.save(str(output))


def main(argv: Sequence[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT, help="caminho do PPTX a gerar")
    parser.add_argument("--check", action="store_true", help="validar o PPTX existente sem gerar outro")
    args = parser.parse_args(argv)
    output = args.output if args.output.is_absolute() else ROOT / args.output

    if args.check:
        errors = validate_presentation(output)
        if errors:
            for error in errors:
                print(f"ERROR: {error}")
            return 1
        print(f"OK: apresentação TCC válida ({len(EXPECTED_TITLES)} slides): {output}")
        return 0

    generate(output)
    errors = validate_presentation(output)
    if errors:
        for error in errors:
            print(f"ERROR: {error}")
        return 1
    print(f"OK: apresentação TCC gerada e validada: {output}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
